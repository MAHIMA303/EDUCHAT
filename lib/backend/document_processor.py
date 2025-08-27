"""
Document processing service for AI Tutor system
Handles various file formats and prepares documents for Haystack
"""

import os
from typing import List, Dict, Any
from haystack.schema import Document
from haystack.nodes import PreProcessor
import logging

# Document processing imports
try:
    from docx import Document as DocxDocument
    from PyPDF2 import PdfReader
    from pptx import Presentation
except ImportError:
    logging.warning("Some document processing libraries not available")

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self):
        self.preprocessor = PreProcessor(
            clean_empty_lines=True,
            clean_whitespace=True,
            clean_header_footer=True,
            split_by="word",
            split_length=200,
            split_overlap=20
        )
    
    def process_text_file(self, file_path: str, subject: str = "General") -> List[Document]:
        """Process a text file and return Haystack documents"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            # Create document with metadata
            doc = Document(
                content=content,
                meta={
                    "source": file_path,
                    "subject": subject,
                    "type": "text",
                    "filename": os.path.basename(file_path)
                }
            )
            
            # Preprocess the document
            processed_docs = self.preprocessor.process([doc])
            return processed_docs
            
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {e}")
            return []
    
    def process_pdf_file(self, file_path: str, subject: str = "General") -> List[Document]:
        """Process a PDF file and return Haystack documents"""
        try:
            reader = PdfReader(file_path)
            content = ""
            
            for page in reader.pages:
                content += page.extract_text() + "\n"
            
            # Create document with metadata
            doc = Document(
                content=content,
                meta={
                    "source": file_path,
                    "subject": subject,
                    "type": "pdf",
                    "filename": os.path.basename(file_path),
                    "pages": len(reader.pages)
                }
            )
            
            # Preprocess the document
            processed_docs = self.preprocessor.process([doc])
            return processed_docs
            
        except Exception as e:
            logger.error(f"Error processing PDF file {file_path}: {e}")
            return []
    
    def process_docx_file(self, file_path: str, subject: str = "General") -> List[Document]:
        """Process a Word document and return Haystack documents"""
        try:
            doc = DocxDocument(file_path)
            content = ""
            
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
            
            # Create document with metadata
            haystack_doc = Document(
                content=content,
                meta={
                    "source": file_path,
                    "subject": subject,
                    "type": "docx",
                    "filename": os.path.basename(file_path)
                }
            )
            
            # Preprocess the document
            processed_docs = self.preprocessor.process([haystack_doc])
            return processed_docs
            
        except Exception as e:
            logger.error(f"Error processing DOCX file {file_path}: {e}")
            return []
    
    def process_pptx_file(self, file_path: str, subject: str = "General") -> List[Document]:
        """Process a PowerPoint presentation and return Haystack documents"""
        try:
            prs = Presentation(file_path)
            content = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            
            # Create document with metadata
            doc = Document(
                content=content,
                meta={
                    "source": file_path,
                    "subject": subject,
                    "type": "pptx",
                    "filename": os.path.basename(file_path),
                    "slides": len(prs.slides)
                }
            )
            
            # Preprocess the document
            processed_docs = self.preprocessor.process([doc])
            return processed_docs
            
        except Exception as e:
            logger.error(f"Error processing PPTX file {file_path}: {e}")
            return []
    
    def process_file(self, file_path: str, subject: str = "General") -> List[Document]:
        """Process any supported file type"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.txt':
            return self.process_text_file(file_path, subject)
        elif file_extension == '.pdf':
            return self.process_pdf_file(file_path, subject)
        elif file_extension == '.docx':
            return self.process_docx_file(file_path, subject)
        elif file_extension == '.pptx':
            return self.process_pptx_file(file_path, subject)
        else:
            logger.warning(f"Unsupported file type: {file_extension}")
            return []
    
    def process_directory(self, directory_path: str, subject: str = "General") -> List[Document]:
        """Process all supported files in a directory"""
        all_documents = []
        
        try:
            for filename in os.listdir(directory_path):
                file_path = os.path.join(directory_path, filename)
                if os.path.isfile(file_path):
                    documents = self.process_file(file_path, subject)
                    all_documents.extend(documents)
            
            logger.info(f"Processed {len(all_documents)} documents from directory")
            return all_documents
            
        except Exception as e:
            logger.error(f"Error processing directory {directory_path}: {e}")
            return []
    
    def create_sample_educational_content(self) -> List[Document]:
        """Create sample educational content for testing"""
        sample_content = [
            {
                "content": """
                Algebra is a branch of mathematics that deals with symbols and the rules for manipulating these symbols. 
                In elementary algebra, those symbols (today written as Latin and Greek letters) represent quantities without 
                fixed values, known as variables. The rules of algebra are used to solve equations and find the values 
                of these variables.
                
                Key concepts in algebra include:
                - Variables and constants
                - Expressions and equations
                - Linear equations
                - Quadratic equations
                - Systems of equations
                """,
                "subject": "Mathematics",
                "topic": "Algebra Basics"
            },
            {
                "content": """
                Physics is the natural science that studies matter, its motion and behavior through space and time, 
                and the related entities of energy and force. Physics is one of the most fundamental scientific 
                disciplines, and its main goal is to understand how the universe behaves.
                
                Main branches of physics:
                - Classical mechanics
                - Thermodynamics
                - Electromagnetism
                - Quantum mechanics
                - Relativity
                """,
                "subject": "Physics",
                "topic": "Introduction to Physics"
            },
            {
                "content": """
                Chemistry is the scientific discipline involved with elements and compounds composed of atoms, 
                molecules and ions: their composition, structure, properties, behavior and the changes they 
                undergo during a reaction with other substances.
                
                Core areas of chemistry:
                - Organic chemistry
                - Inorganic chemistry
                - Physical chemistry
                - Biochemistry
                - Analytical chemistry
                """,
                "subject": "Chemistry",
                "topic": "Chemistry Fundamentals"
            }
        ]
        
        documents = []
        for content in sample_content:
            doc = Document(
                content=content["content"],
                meta={
                    "subject": content["subject"],
                    "topic": content["topic"],
                    "type": "sample",
                    "source": "sample_content"
                }
            )
            documents.append(doc)
        
        return documents

# Global instance
document_processor = DocumentProcessor()
